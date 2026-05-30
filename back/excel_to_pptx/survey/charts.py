import pandas as pd
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from back.excel_to_pptx.constants import Constants

logger = logging.getLogger(__name__)


def create_chart_slides(prs: Presentation, df: pd.DataFrame, layout):
    """Создает слайды с круговыми диаграммами для всех подходящих вопросов"""
    
    from .analyzer import SurveyAnalyzer
    
    analyzer = SurveyAnalyzer(df)
    analyzer.group_age_column()
    
    slide_width_inches = prs.slide_width.inches
    
    for col in df.columns:
        if analyzer.is_technical_or_excluded_col(col):
            continue

        valid_series = df[col].dropna().astype(str).str.strip()
        unique_answers = valid_series.unique()

        if len(unique_answers) > 8 or len(unique_answers) <= 1:
            continue

        val_counts = valid_series.value_counts()
        if val_counts.empty:
            continue

        logger.info(f"Создание слайда для: {col[:50]}...")

        slide = prs.slides.add_slide(layout)

        for shape in list(slide.shapes):
            if shape.is_placeholder:
                sp = shape._element
                sp.getparent().remove(sp)

        header_line = slide.shapes.add_shape(1, *Constants.HEADER_LINE_GEOMETRY)
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = Constants.COLOR_ACCENT
        header_line.line.fill.background()

        tx_title = slide.shapes.add_textbox(*Constants.HEADER_TEXT_BOX_GEOMETRY)
        tf_title = tx_title.text_frame
        tf_title.word_wrap = True
        p_t = tf_title.paragraphs[0]
        p_t.text = col
        p_t.font.name = Constants.FONT_NAME
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = Constants.COLOR_PRIMARY

        total_responses = val_counts.sum()
        percentages = val_counts / total_responses

        chart_data = CategoryChartData()
        categories = []
        for cat in val_counts.index:
            cat_str = str(cat)
            if len(cat_str) > 40:
                cat_str = cat_str[:37] + "..."
            categories.append(cat_str)

        chart_data.categories = categories
        chart_data.add_series('Ответы', tuple(percentages.values))

        chart_w = Inches(8.0)
        chart_h = Inches(5.5)
        chart_w_inches = 8.0
        chart_l = Inches((slide_width_inches - chart_w_inches) / 2)
        chart_t = Inches(2.0)

        try:
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE,
                chart_l, chart_t, chart_w, chart_h,
                chart_data
            )
            chart = chart_shape.chart

            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.font.name = Constants.FONT_NAME
            chart.legend.font.size = Pt(18)
            chart.legend.include_in_layout = False

            plot = chart.plots[0]
            series = plot.series[0]

            for idx, point in enumerate(series.points):
                color_idx = idx % len(Constants.SURVEY_COLORS)
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = Constants.SURVEY_COLORS[color_idx]

            plot.has_data_labels = True
            data_labels = plot.data_labels

            data_labels.number_format = '0%'
            data_labels.position = XL_LABEL_POSITION.BEST_FIT
            data_labels.font.name = Constants.FONT_NAME
            data_labels.font.size = Pt(16)

            if len(val_counts) > 6:
                data_labels.font.size = Pt(13)
            elif len(val_counts) > 8:
                data_labels.font.size = Pt(11)

            logger.info(f"Диаграмма для '{col[:30]}...' создана")

        except Exception as e:
            logger.error(f"Ошибка при создании диаграммы для {col}: {e}")
            continue