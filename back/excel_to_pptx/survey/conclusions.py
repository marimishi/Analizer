from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from back.excel_to_pptx.constants import Constants
import logging

logger = logging.getLogger(__name__)


def create_conclusion_slides(prs: Presentation, user_choices: dict, layout):
    """
    Создает слайды с выводами на основе выбора пользователя
    """
    
    categories_to_render = [
        {"title": "Позитивные аспекты работы организации", 
         "key": "Позитивный аспект", 
         "accent_color": Constants.COLOR_THIRDLY},
        {"title": "Проблемные области и зоны развития", 
         "key": "Проблемная область", 
         "accent_color": Constants.COLOR_ERROR}
    ]
    
    for cat in categories_to_render:
        questions = user_choices.get(cat["key"], [])
        if not questions:
            logger.info(f"Нет вопросов для категории: {cat['key']}")
            continue

        logger.info(f"Создание слайда для категории: {cat['key']}, вопросов: {len(questions)}")
        
        slide = prs.slides.add_slide(layout)

        # Очистка плейсхолдеров
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                sp = shape._element
                sp.getparent().remove(sp)

        header_line = slide.shapes.add_shape(1, *Constants.HEADER_LINE_GEOMETRY)
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = cat["accent_color"]
        header_line.line.fill.background()

        tx_title = slide.shapes.add_textbox(*Constants.HEADER_TEXT_BOX_GEOMETRY)
        p_t = tx_title.text_frame.paragraphs[0]
        p_t.text = cat["title"]
        p_t.font.name = Constants.FONT_NAME
        p_t.font.size = Constants.HEADER_FONT_SIZE
        p_t.font.bold = True
        p_t.font.color.rgb = Constants.COLOR_PRIMARY

        text_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), 
            Inches(11.7), Inches(5.0)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for i, q_text in enumerate(questions):
            p_num = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p_num.text = f"{i+1}."
            p_num.font.name = Constants.FONT_NAME
            p_num.font.size = Pt(18)
            p_num.font.bold = True
            p_num.font.color.rgb = Constants.COLOR_ACCENT
            p_num.space_before = Pt(12) if i > 0 else Pt(0)
            
            p_text = tf.add_paragraph()
            clean_q = q_text.replace('\n', ' ').strip()
            p_text.text = clean_q
            p_text.font.name = Constants.FONT_NAME
            p_text.font.size = Pt(16)
            p_text.font.color.rgb = Constants.COLOR_PRIMARY
            p_text.space_before = Pt(4)
            p_text.line_spacing = 1.3