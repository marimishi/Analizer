import pandas as pd
import re
import os
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION
from back.excel_to_pptx.constants import Constants
from back.excel_to_pptx.title import create_title_slide
from back.yandex_forms.excel_to_pandas import ExcelToPandasProcessor

# Настройка логирования
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

def is_technical_or_excluded_col(col_name: str) -> bool:
    col_clean = col_name.strip().lower()
    exclude_patterns = [
        r'\bid\b', r'идентификатор', r'фио', r'имя', r'фамилия', 
        r'дата', r'время', r'ваши предложения'
    ]
    for pattern in exclude_patterns:
        if re.search(pattern, col_clean):
            return True
    return False

def is_demographic_col(col_name: str) -> bool:
    col_clean = col_name.strip().lower()
    return 'пол' in col_clean or 'возраст' in col_clean

def calculate_column_score(series: pd.Series) -> float:
    """
    Возвращает процент позитивных ответов от общего числа.
    Чем ниже процент, тем сильнее это 'Проблемная область'.
    """
    valid_data = series.dropna().astype(str).str.lower().str.strip()
    if valid_data.empty:
        return 0.0
        
    total_count = len(valid_data)
    pos_patterns = [r'\bда\b', r'полностью', r'отлично', r'хорошо', r'удовлетворен', r'легко', r'рекомендую']
    
    pos_clicks = 0
    for val in valid_data:
        if any(p in val for p in pos_patterns) and "не " not in val and "плохо" not in val:
            pos_clicks += 1
            
    return pos_clicks / total_count

def survey_base_diagram(df: pd.DataFrame, prs: Presentation, layout=None, file_path: str = None, title_text: str = None) -> list:
    """
    Анализирует датафрейм, создает титульный слайд, затем ранжирует 
    вопросы по удовлетворенности и генерирует круговые диаграммы по центру слайда.
    """
    logger.info("Начало создания презентации с диаграммами")
    
    # Если UI передает пустую строку или None, берем имя файла и чистим его методом Яндекса
    final_title = title_text.strip() if title_text else ""
    if not final_title and file_path:
        final_title = ExcelToPandasProcessor.clean_filename_from_date(file_path)
    elif not final_title:
        final_title = "Отчет по результатам анкетирования"

    create_title_slide(prs, final_title, prs.slide_layouts[0])
    
    raw_computed_cols = []
    
    for col in df.columns:
        if is_technical_or_excluded_col(col):
            continue
            
        valid_series = df[col].dropna().astype(str).str.strip()
        unique_answers = valid_series.unique()
        
        if len(unique_answers) > 8 or len(unique_answers) <= 1:
            continue
            
        is_demographic = is_demographic_col(col)
        score = 0.0 if is_demographic else calculate_column_score(df[col])
        
        raw_computed_cols.append({
            "column_name": col,
            "score": score,
            "is_demographic": is_demographic,
            "valid_series": valid_series
        })
        
    if not raw_computed_cols:
        return []

    scoring_questions = [q for q in raw_computed_cols if not q["is_demographic"]]
    scoring_questions.sort(key=lambda x: x["score"])
    
    total_valid_questions = len(scoring_questions)
    
    aspect_predictions = {}
    for index, item in enumerate(scoring_questions):
        if index < 3:
            aspect_predictions[item["column_name"]] = "Проблемная область"
        elif index >= (total_valid_questions - 3):
            aspect_predictions[item["column_name"]] = "Позитивный аспект"
        else:
            aspect_predictions[item["column_name"]] = "Не включать в выводы"
            
    for item in raw_computed_cols:
        if item["is_demographic"]:
            aspect_predictions[item["column_name"]] = "Не включать в выводы"
    
    generated_charts_config = []
    slide_width = prs.slide_width
    
    try:
        blank_layout = prs.slide_layouts[6]
    except:
        blank_layout = prs.slide_layouts[1]
    
    for item in raw_computed_cols:
        col = item["column_name"]
        valid_series = item["valid_series"]
        predicted_aspect = aspect_predictions[col]
        
        val_counts = valid_series.value_counts()
        if val_counts.empty:
            continue
            
        slide = prs.slides.add_slide(blank_layout)
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape.text_frame.clear()
        
        header_line = slide.shapes.add_shape(1, *Constants.HEADER_LINE_GEOMETRY)
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = Constants.COLOR_ACCENT
        header_line.line.fill.background()

        tx_title = slide.shapes.add_textbox(*Constants.HEADER_TEXT_BOX_GEOMETRY)
        tf_title = tx_title.text_frame
        tf_title.word_wrap = True
        p_t = tf_title.paragraphs[0]
        p_t.text = col
        p_t.font.name = Constants.FONT_NAME
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = Constants.COLOR_PRIMARY
        
        total_responses = val_counts.sum()
        percentages = val_counts / total_responses
        
        chart_data = CategoryChartData()
        categories = []
        for cat in val_counts.index:
            cat_str = str(cat)
            if len(cat_str) > 40:
                cat_str = cat_str[:37] + "..."
            categories.append(cat_str)
        
        chart_data.categories = categories
        chart_data.add_series('Ответы', tuple(percentages.values))
        
        chart_w = Inches(5.8)
        chart_h = Inches(4.8)
        chart_l = (slide_width - chart_w) / 2
        chart_t = Inches(2.3)
        
        try:
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE,  # Меняем на PIE
                chart_l, chart_t, chart_w, chart_h, 
                chart_data
            )
            chart = chart_shape.chart
            
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.font.name = Constants.FONT_NAME
            chart.legend.font.size = Pt(10)
            chart.legend.include_in_layout = False
            
            plot = chart.plots[0]
            series = plot.series[0]
            
            for idx, point in enumerate(series.points):
                color_idx = idx % len(Constants.SURVEY_COLORS)
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = Constants.SURVEY_COLORS[color_idx]
            
            plot.has_data_labels = True
            data_labels = plot.data_labels
            
            data_labels.number_format = '0%'
            
            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            
            data_labels.font.name = Constants.FONT_NAME
            data_labels.font.size = Pt(11)
            data_labels.font.bold = True
            
            try:
                data_labels.show_percentage = True
                data_labels.show_value = False
                data_labels.show_category_name = False
            except:
                pass
            
            if len(val_counts) > 6:
                data_labels.font.size = Pt(9)
            elif len(val_counts) > 8:
                data_labels.font.size = Pt(8)
            
            logger.info(f"Круговая диаграмма для '{col}' создана ({len(val_counts)} секторов)")
            
        except Exception as e:
            logger.error(f"Ошибка при создании диаграммы для {col}: {e}")
            continue
        
        generated_charts_config.append({
            "column_name": col,
            "predicted_aspect": predicted_aspect,
            "answers_preview": dict(val_counts)
        })
        
    return generated_charts_config

def create_presentation(excel_file_path: str, output_pptx_path: str, title_text: str = None) -> bool:
    """
    Создает презентацию из Excel файла
    """
    try:
        logger.info(f"Загрузка Excel файла: {excel_file_path}")
        processor = ExcelToPandasProcessor()
        df = processor.load_excel_to_dataframe(excel_file_path)
        logger.info(f"Загружено {len(df)} строк, {len(df.columns)} колонок")
        
        prs = Presentation()
        
        charts = survey_base_diagram(df, prs, file_path=excel_file_path, title_text=title_text)
        logger.info(f"Создано {len(charts)} слайдов с круговыми диаграммами")
        
        output_dir = os.path.dirname(output_pptx_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        if os.path.exists(output_pptx_path):
            try:
                os.remove(output_pptx_path)
            except PermissionError:
                logger.error(f"Файл {output_pptx_path} открыт в другой программе")
                return False
        
        prs.save(output_pptx_path)
        
        if os.path.exists(output_pptx_path) and os.path.getsize(output_pptx_path) > 1000:
            return True
        else:
            return False
            
    except Exception as e:
        logger.error(f"Ошибка при создании презентации: {e}", exc_info=True)
        return False

if __name__ == "__main__":
    success = create_presentation(
        excel_file_path="path/to/your/excel.xlsx",
        output_pptx_path="output/presentation.pptx",
        title_text="Отчет по удовлетворенности пациентов"
    )
    
    if success:
        print("Презентация успешно создана!")
    else:
        print("Ошибка при создании презентации")