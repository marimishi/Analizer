import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from back.excel_to_pptx.constants import Constants


def create_common_results_slide(prs: Presentation, df: pd.DataFrame, layout):
    slide = prs.slides.add_slide(layout)
    
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)

    total_participants = len(df)
    
    text_block_height = Inches(3.5)
    top = (prs.slide_height - text_block_height) / 2
    
    tx_center = slide.shapes.add_textbox(
        Inches(0), top, 
        Inches(13.333), text_block_height
    )
    
    tf = tx_center.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = f"{total_participants}"
    p1.font.name = Constants.FONT_NAME
    p1.font.size = Pt(166)
    p1.font.bold = True
    p1.font.color.rgb = Constants.COLOR_ACCENT
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "ВСЕГО УЧАСТНИКОВ ОПРОСА"
    p2.font.name = Constants.FONT_NAME
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = Constants.COLOR_PRIMARY
    p2.space_before = Pt(20)
    p2.alignment = PP_ALIGN.CENTER