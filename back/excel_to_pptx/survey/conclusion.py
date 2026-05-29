import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from back.excel_to_pptx.constants import Constants

def create_survey_conclusions(prs: Presentation, user_choices: dict, layout):
    """
    user_choices: словарь вида:
    {
        "Позитивный аспект": ["Вопрос 1", "Вопрос 3"],
        "Проблемная область": ["Вопрос 2", "Вопрос 5"],
        "Не включать в выводы": ["Вопрос 4"]
    }
    """
    
    categories_to_render = [
        {"title": "Позитивные аспекты работы организации", "key": "Позитивный аспект", "accent_color": Constants.COLOR_THIRDLY},
        {"title": "Проблемные области и зоны развития", "key": "Проблемная область", "accent_color": Constants.COLOR_ERROR}
    ]
    
    for cat in categories_to_render:
        questions = user_choices.get(cat["key"], [])
        if not questions:
            continue # Если в категорию ничего не выбрали, пропускаем слайд или делаем пустым
            
        slide = prs.slides.add_slide(layout)
        
        # Очистка плейсхолдеров
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                sp = shape._element
                sp.getparent().remove(sp)
                
        # Боковая цветная линия-маркер под тему
        header_line = slide.shapes.add_shape(1, *Constants.HEADER_LINE_GEOMETRY)
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = cat["accent_color"]
        header_line.line.fill.background()

        # Заголовок слайда вывода
        tx_title = slide.shapes.add_textbox(*Constants.HEADER_TEXT_BOX_GEOMETRY)
        p_t = tx_title.text_frame.paragraphs[0]
        p_t.text = cat["title"]
        p_t.font.name = Constants.FONT_NAME
        p_t.font.size = Constants.HEADER_FONT_SIZE
        p_t.font.bold = True
        p_t.font.color.rgb = Constants.COLOR_PRIMARY
        
        # Контейнер списков выводов (Плиточный или Списковый вывод на левую и правую часть экрана)
        # Для гибкости распределим до 4-6 пунктов в виде красивых карточек-текстбоксов
        for i, q_text in enumerate(questions[:4]): # Лимит 4 карточки на слайд для сохранения дизайна
            # Рассчитываем координаты сетки 2х2
            is_right = i % 2 == 1
            is_bottom = i >= 2
            
            x = Constants.TILE_X_RIGHT if is_right else Constants.TILE_X_LEFT
            y = Constants.TILE_Y_BOTTOM if is_bottom else Constants.TILE_Y_TOP
            
            # Фоновая плашка пункта выводов
            tile_shape = slide.shapes.add_shape(1, x, y, Constants.TILE_W, Constants.TILE_H)
            tile_shape.fill.solid()
            tile_shape.fill.fore_color.rgb = Constants.COLOR_TILE_BG
            tile_shape.line.color.rgb = Constants.COLOR_SECONDARY
            tile_shape.line.width = Constants.TILE_LINE_WIDTH
            if hasattr(tile_shape, "shadow") and tile_shape.shadow:
                tile_shape.shadow.inherit = False
                
            # Текст внутри плашки
            tx_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Constants.TILE_W - Inches(0.6), Constants.TILE_H - Inches(0.4))
            tf = tx_box.text_frame
            tf.word_wrap = True
            
            p_tag = tf.paragraphs[0]
            p_tag.text = f"КЛЮЧЕВОЙ ПУНКТ #{i+1}"
            p_tag.font.name = Constants.FONT_NAME
            p_tag.font.size = Pt(12)
            p_tag.font.bold = True
            p_tag.font.color.rgb = Constants.COLOR_MUTED_TAG
            
            p_main = tf.add_paragraph()
            # Форматируем текст вывода
            clean_q = q_text.replace('\n', ' ').strip()
            p_main.text = clean_q if len(clean_q) <= 120 else clean_q[:117] + "..."
            p_main.font.name = Constants.FONT_NAME
            p_main.font.size = Pt(16)
            p_main.font.bold = True
            p_main.font.color.rgb = Constants.COLOR_PRIMARY
            p_main.space_before = Pt(6)
            p_main.line_spacing = 1.2