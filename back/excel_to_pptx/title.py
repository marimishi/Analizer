from pptx import Presentation
from pptx.util import Inches, Pt
from back.excel_to_pptx.constants import Constants

def create_title_slide(prs: Presentation, title_text: str, layout):
    slide = prs.slides.add_slide(layout)
    
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)

    FONT_NAME = Constants.FONT_NAME
    COLOR_PRIMARY = Constants.COLOR_PRIMARY
    COLOR_SECONDARY = Constants.COLOR_SECONDARY
    COLOR_ACCENT = Constants.COLOR_ACCENT
    SUBTITLE_TEXT = 'ГАУЗ ТО "МКДЦ"'

    line = slide.shapes.add_shape(
        1, Inches(1.0), Inches(2.6), Inches(0.04), Inches(2.2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    left, top = Inches(1.3), Inches(2.5)
    width, height = Inches(10), Inches(3)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = 0

    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_NAME
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.line_spacing = 1.1

    p2 = tf.add_paragraph()
    p2.text = SUBTITLE_TEXT
    p2.font.name = FONT_NAME
    p2.font.size = Pt(22)
    p2.font.bold = False
    p2.font.color.rgb = COLOR_SECONDARY
    p2.space_before = Pt(12)