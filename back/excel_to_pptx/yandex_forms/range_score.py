import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from back.excel_to_pptx.constants import Constants

def create_range_score_slide(prs: Presentation, df: pd.DataFrame, layout):
    slide = prs.slides.add_slide(layout)
    
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)

    header_line = slide.shapes.add_shape(1, *Constants.HEADER_LINE_GEOMETRY)
    header_line.fill.solid()
    header_line.fill.fore_color.rgb = Constants.COLOR_ACCENT
    header_line.line.fill.background()

    tx_title = slide.shapes.add_textbox(*Constants.HEADER_TEXT_BOX_GEOMETRY)
    p_t = tx_title.text_frame.paragraphs[0]
    p_t.text = "Распределение участников по набранным баллам"
    p_t.font.name = Constants.FONT_NAME
    p_t.font.size = Constants.HEADER_FONT_SIZE
    p_t.font.bold = True
    p_t.font.color.rgb = Constants.COLOR_PRIMARY

    score_cols = [c for c in df.columns if "/ Баллы" in c]
    user_scores = df[score_cols].apply(pd.to_numeric, errors='coerce').fillna(0).sum(axis=1).round().astype(int)
    score_counts = user_scores.value_counts().sort_index()
    
    chart_data = CategoryChartData()
    chart_data.categories = [f"{b} б." for b in score_counts.index]
    chart_data.add_series('Количество человек', tuple(int(v) for v in score_counts.values))
    
    chart_x = Inches(0.8)
    chart_y = Inches(1.8)
    chart_width = Inches(11.7)
    chart_height = Inches(4.8)
    
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, 
        chart_x, chart_y, chart_width, chart_height, 
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    
    series = chart.plots[0].series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = Constants.COLOR_ACCENT
    
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.name = Constants.FONT_NAME
    data_labels.font.size = Constants.CHART_LABEL_SIZE
    data_labels.font.bold = True
    data_labels.font.color.rgb = Constants.COLOR_PRIMARY
    
    data_labels.show_percentage = False
    data_labels.show_value = True
    data_labels.number_format = '0'