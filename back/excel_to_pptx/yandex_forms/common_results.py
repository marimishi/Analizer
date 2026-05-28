import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from back.excel_to_pptx.constants import Constants

def create_common_results_slide(prs: Presentation, df: pd.DataFrame, layout):
    slide = prs.slides.add_slide(layout)
    
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)

    header_line = slide.shapes.add_shape(1, *Constants.HEADER_LINE_GEOMETRY)
    header_line.fill.solid()
    header_line.fill.fore_color.rgb = Constants.COLOR_ACCENT
    header_line.line.fill.background()

    tx_title = slide.shapes.add_textbox(*Constants.HEADER_TEXT_BOX_GEOMETRY)
    p_t = tx_title.text_frame.paragraphs[0]
    p_t.text = "Общие результаты тестирования"
    p_t.font.name = Constants.FONT_NAME
    p_t.font.size = Constants.HEADER_FONT_SIZE
    p_t.font.bold = True
    p_t.font.color.rgb = Constants.COLOR_PRIMARY

    total_participants = len(df)
    score_cols = [c for c in df.columns if "/ Баллы" in c]
    total_questions = len(score_cols)
    user_scores = df[score_cols].apply(pd.to_numeric, errors='coerce').fillna(0).sum(axis=1)
    max_possible_score = total_questions
    mean_score = user_scores.mean()
    success_rate = (mean_score / max_possible_score) * 100 if max_possible_score > 0 else 0

    metrics_dataset = [
        {"value": f"{total_participants}", "label": "участников успешно прошли тестирование"},
        {"label": f"количество вопросов", "value": f"{total_questions}"},
        {"label": f"со средним баллом", "value": f"{mean_score:.2f} из {max_possible_score}"},
        {"value": f"{success_rate:.0f}%", "label": "общая успешность выполнения заданий", "is_accent": True}
    ]

    for i, data in enumerate(metrics_dataset):
        current_y = Constants.METRIC_START_Y + i * (Constants.METRIC_HEIGHT + Constants.METRIC_GAP)
        
        tx_metric = slide.shapes.add_textbox(
            Constants.METRIC_START_X, current_y, 
            Constants.METRIC_WIDTH, Constants.METRIC_HEIGHT
        )
        tf = tx_metric.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p1 = tf.paragraphs[0]
        p1.text = data["value"]
        p1.font.name = Constants.FONT_NAME
        p1.font.size = Constants.METRIC_VALUE_SIZE
        p1.font.bold = True
        p1.font.color.rgb = Constants.COLOR_ACCENT if data.get("is_accent") else Constants.COLOR_PRIMARY
        
        p2 = tf.add_paragraph()
        p2.text = data["label"]
        p2.font.name = Constants.FONT_NAME
        p2.font.size = Constants.METRIC_LABEL_SIZE
        p2.font.color.rgb = Constants.COLOR_SECONDARY
        p2.space_before = Pt(2)

    sep_line = slide.shapes.add_shape(
        1, 
        Inches(6.3),
        Inches(1.5),
        Inches(0.01),
        Inches(5.5)
    )
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = Constants.COLOR_SECONDARY
    sep_line.line.fill.background()

    chart_data = CategoryChartData()
    chart_data.categories = ['Успешно', 'Неуспешно']
    
    success_fraction = success_rate / 100
    chart_data.add_series('Результат', (success_fraction, 1.0 - success_fraction))
    
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, *Constants.CHART_GEOMETRY, chart_data)
    chart = chart_shape.chart
    
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.name = Constants.FONT_NAME
    chart.legend.font.size = Constants.CHART_LEGEND_SIZE
    
    points = chart.plots[0].series[0].points
    points[0].format.fill.solid()
    points[0].format.fill.fore_color.rgb = Constants.COLOR_ACCENT
    points[1].format.fill.solid()
    points[1].format.fill.fore_color.rgb = Constants.COLOR_THIRDLY

    plot = chart.plots[0]
    plot.has_data_labels = True
    
    try:
        plot.hole_size = 70
    except:
        pass

    data_labels = plot.data_labels
    data_labels.font.name = Constants.FONT_NAME
    data_labels.font.size = Constants.CHART_LABEL_SIZE
    data_labels.font.bold = True
    data_labels.font.color.rgb = Constants.COLOR_PRIMARY
    
    data_labels.show_percentage = False
    data_labels.show_value = True 
    data_labels.number_format = '0%'