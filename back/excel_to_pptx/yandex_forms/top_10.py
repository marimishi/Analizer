import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from back.excel_to_pptx.constants import Constants

def _base_create_chart_slide(
    prs: Presentation, 
    layout, 
    slide_title: str, 
    categories: list, 
    values: list, 
    series_name: str,
    number_format: str = '0',
    max_value_scale: float = None
):
    """Базовая внутренняя функция для создания слайда с кастомной столбчатой диаграммой."""
    slide = prs.slides.add_slide(layout)
    
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)

    header_line = slide.shapes.add_shape(1, *Constants.HEADER_LINE_GEOMETRY)
    header_line.fill.solid()
    header_line.fill.fore_color.rgb = Constants.COLOR_ACCENT
    header_line.line.fill.background()
    header_line.shadow.inherit = False

    tx_title = slide.shapes.add_textbox(*Constants.HEADER_TEXT_BOX_GEOMETRY)
    p_t = tx_title.text_frame.paragraphs[0]
    p_t.text = slide_title
    p_t.font.name = Constants.FONT_NAME
    p_t.font.size = Constants.HEADER_FONT_SIZE
    p_t.font.bold = True
    p_t.font.color.rgb = Constants.COLOR_PRIMARY

    chart_data = CategoryChartData()
    chart_data.categories = [str(name)[:20] + "..." if len(str(name)) > 20 else str(name) for name in categories]
    chart_data.add_series(series_name, tuple(values))
    
    chart_x, chart_y = Inches(0.8), Inches(1.8)
    chart_width, chart_height = Inches(11.7), Inches(4.8)
    
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, 
        chart_x, chart_y, chart_width, chart_height, 
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    
    series = chart.plots[0].series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = Constants.COLOR_ACCENT
    
    category_axis = chart.category_axis
    category_axis.tick_labels.font.name = Constants.FONT_NAME
    category_axis.tick_labels.font.size = Constants.CHART_LABEL_SIZE
    category_axis.tick_labels.font.color.rgb = Constants.COLOR_PRIMARY

    value_axis = chart.value_axis
    if max_value_scale is not None:
        value_axis.maximum_scale = max_value_scale
    value_axis.tick_labels.font.name = Constants.FONT_NAME
    value_axis.tick_labels.font.size = Constants.CHART_LABEL_SIZE
    value_axis.tick_labels.font.color.rgb = Constants.COLOR_SECONDARY
    value_axis.number_format = '0' # Сама ось Y всегда отображает просто число
    
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.name = Constants.FONT_NAME
    data_labels.font.size = Constants.CHART_LABEL_SIZE
    data_labels.font.bold = True
    data_labels.font.color.rgb = Constants.COLOR_PRIMARY
    
    data_labels.show_percentage = False
    data_labels.show_value = True
    data_labels.number_format = number_format


def create_top_10_slide(prs: Presentation, df: pd.DataFrame, layout):
    dep_col = "Укажите подразделение/отделение"
    if dep_col not in df.columns:
        return
        
    top10 = df[dep_col].value_counts()
    
    _base_create_chart_slide(
        prs=prs,
        layout=layout,
        slide_title="Топ подразделений по количеству участников",
        categories=list(top10.index),
        values=[int(v) for v in top10.values],
        series_name='Сотрудников',
        number_format='0'
    )


def create_top_10_success_slide(prs: Presentation, df: pd.DataFrame, layout):
    dep_col = "Укажите подразделение/отделение"
    if dep_col not in df.columns:
        return

    score_cols = [c for c in df.columns if "/ Баллы" in c]
    df_temp = df.copy()
    df_temp['Total_Score'] = df_temp[score_cols].apply(pd.to_numeric, errors='coerce').fillna(0).sum(axis=1)
    max_score = len(score_cols)
    
    summary = df_temp.groupby(dep_col).agg(
        Средний_балл=('Total_Score', 'mean')
    ).reset_index()
    
    summary['Успешность'] = (summary['Средний_балл'] / max_score * 100).round(1) if max_score > 0 else 0
    top10_success = summary.sort_values(by='Успешность', ascending=False)
    
    _base_create_chart_slide(
        prs=prs,
        layout=layout,
        slide_title="Топ подразделений по успешности тестирования",
        categories=list(top10_success[dep_col]),
        values=[float(v) for v in top10_success['Успешность']],
        series_name='Успешность',
        number_format='0"%"',
        max_value_scale=100
    )