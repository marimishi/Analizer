import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from back.excel_to_pptx.constants import Constants

def create_department_slides(prs: Presentation, df: pd.DataFrame, layout):
    dep_col = "Укажите подразделение/отделение"
    if dep_col not in df.columns:
        return
        
    score_cols = [c for c in df.columns if "/ Баллы" in c]
    df['Total_Score'] = df[score_cols].apply(pd.to_numeric, errors='coerce').fillna(0).sum(axis=1)
    max_score = len(score_cols)
    
    summary = df.groupby(dep_col).agg(
        Участников=('Total_Score', 'count'),
        Средний_балл=('Total_Score', 'mean')
    ).reset_index()
    
    summary['Успешность_число'] = (summary['Средний_балл'] / max_score * 100).round(1) if max_score > 0 else 0
    summary['Успешность'] = summary['Успешность_число'].astype(str) + "%"
    summary['Middle_Score'] = summary['Средний_балл'].round(2)
    
    chunk_size = 7
    chunks = [summary[i:i + chunk_size] for i in range(0, summary.shape[0], chunk_size)]
    
    for idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(layout)
        
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                sp = shape._element
                sp.getparent().remove(sp)

        header_line = slide.shapes.add_shape(1, *Constants.HEADER_LINE_GEOMETRY)
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = Constants.COLOR_ACCENT
        header_line.line.fill.background()

        tx_title = slide.shapes.add_textbox(*Constants.HEADER_TEXT_BOX_GEOMETRY)
        p_t = tx_title.text_frame.paragraphs[0]
        page_suffix = f" (Часть {idx+1})" if len(chunks) > 1 else ""
        p_t.text = f"Результаты по подразделениям{page_suffix}"
        p_t.font.name = Constants.FONT_NAME
        p_t.font.size = Constants.HEADER_FONT_SIZE
        p_t.font.bold = True
        p_t.font.color.rgb = Constants.COLOR_PRIMARY
        
        rows = len(chunk) + 1
        cols = 4
        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(11.7)
        row_height_val = 0.6
        height = Inches(row_height_val * rows)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        table.columns[0].width = Inches(5.2)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(2.2)
        table.columns[3].width = Inches(2.3)
        
        headers = ["Подразделение", "Участников", "Средний балл", "Успешность"]
        for col_idx, text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.fill.background()
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.font.name = Constants.FONT_NAME
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = Constants.COLOR_PRIMARY
            if col_idx > 0:
                p.alignment = PP_ALIGN.CENTER
                
        line_y = top + Inches(row_height_val)
        line = slide.shapes.add_shape(1, left, line_y, width, Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = Constants.COLOR_PRIMARY
        line.line.fill.background()
        line.shadow.inherit = False
            
        for row_idx, (_, row) in enumerate(chunk.iterrows(), start=1):
            cell_dep = table.cell(row_idx, 0)
            cell_dep.fill.background()
            cell_dep.text = str(row[dep_col])
            p_dep = cell_dep.text_frame.paragraphs[0]
            p_dep.font.name = Constants.FONT_NAME
            p_dep.font.size = Pt(18)
            p_dep.font.color.rgb = Constants.COLOR_PRIMARY
            
            cell_cnt = table.cell(row_idx, 1)
            cell_cnt.fill.background()
            cell_cnt.text = str(row['Участников'])
            p_cnt = cell_cnt.text_frame.paragraphs[0]
            p_cnt.font.name = Constants.FONT_NAME
            p_cnt.font.size = Pt(18)
            p_cnt.font.color.rgb = Constants.COLOR_PRIMARY
            p_cnt.alignment = PP_ALIGN.CENTER
            
            cell_score = table.cell(row_idx, 2)
            cell_score.fill.background()
            cell_score.text = str(row['Middle_Score'])
            p_score = cell_score.text_frame.paragraphs[0]
            p_score.font.name = Constants.FONT_NAME
            p_score.font.size = Pt(18)
            p_score.font.color.rgb = Constants.COLOR_PRIMARY
            p_score.alignment = PP_ALIGN.CENTER
            
            cell_pct = table.cell(row_idx, 3)
            cell_pct.text = str(row['Успешность'])
            p_pct = cell_pct.text_frame.paragraphs[0]
            p_pct.font.name = Constants.FONT_NAME
            p_pct.font.size = Pt(18)
            p_pct.font.bold = True
            p_pct.alignment = PP_ALIGN.CENTER
            
            val = row['Успешность_число']
            if val >= 80.0:
                bg_color = Constants.COLOR_THIRDLY
            elif val >= 50.0:
                bg_color = Constants.COLOR_WARNING
            else:
                bg_color = Constants.COLOR_ERROR

            cell_pct.fill.solid()
            cell_pct.fill.fore_color.rgb = bg_color
            p_pct.font.color.rgb = Constants.COLOR_PRIMARY
            
            current_line_y = top + Inches(row_height_val * (row_idx + 1))
            row_line = slide.shapes.add_shape(1, left, current_line_y, width, Inches(0.005))
            row_line.fill.solid()
            row_line.fill.fore_color.rgb = Constants.COLOR_SECONDARY
            row_line.line.fill.background()
            row_line.shadow.inherit = False