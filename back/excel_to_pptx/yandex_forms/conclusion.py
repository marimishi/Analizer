import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from back.excel_to_pptx.constants import Constants

def create_conclusion_slide(prs: Presentation, df: pd.DataFrame, layout):
    slide = prs.slides.add_slide(layout)
    
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)

    header_line = slide.shapes.add_shape(1, *Constants.HEADER_LINE_GEOMETRY)
    header_line.fill.solid()
    header_line.fill.fore_color.rgb = Constants.COLOR_ACCENT
    header_line.line.fill.background()

    tx_title = slide.shapes.add_textbox(*Constants.HEADER_TEXT_BOX_GEOMETRY)
    p_t = tx_title.text_frame.paragraphs[0]
    p_t.text = "Заключение и Выводы"
    p_t.font.name = Constants.FONT_NAME
    p_t.font.size = Constants.HEADER_FONT_SIZE
    p_t.font.bold = True
    p_t.font.color.rgb = Constants.COLOR_PRIMARY

    total_participants = len(df)
    score_cols = [c for c in df.columns if "/ Баллы" in c]
    max_score = len(score_cols)
    
    user_scores = df[score_cols].apply(pd.to_numeric, errors='coerce').fillna(0).sum(axis=1)
    mean_score = user_scores.mean()
    success_rate = (mean_score / max_score) * 100 if max_score > 0 else 0
    
    passed_count = sum(user_scores >= (max_score * 0.7))
    passed_pct = (passed_count / total_participants) * 100
    
    status = "Высокий" if success_rate >= 80 else ("Средний" if success_rate >= 50 else "Низкий")
    comparison = "превышает установленный порог" if success_rate >= 70 else "меньше целевого показателя"
    
    q_success = {col: (pd.to_numeric(df[col], errors='coerce').fillna(0) > 0).mean() * 100 for col in score_cols}
    worst_q = min(q_success, key=q_success.get).replace("/ Баллы", "").strip() if q_success else "Нет данных"
    
    best_dep, worst_dep = "Нет данных", "Нет данных"
    best_dep_pct, worst_dep_pct = 0, 0
    
    dep_col = "Укажите подразделение/отделение"
    df_temp = df.copy()
    df_temp['Total_Score'] = user_scores
    
    if dep_col in df_temp.columns:
        dep_stats = df_temp.groupby(dep_col)['Total_Score'].mean() / max_score * 100
        if not dep_stats.empty:
            best_dep = dep_stats.idxmax()
            best_dep_pct = dep_stats.max()
            worst_dep = dep_stats.idxmin()
            worst_dep_pct = dep_stats.min()

    tiles_config = [
        (Constants.TILE_X_LEFT, Constants.TILE_Y_TOP, "РЕЗЮМЕ АНАЛИТИКА", status.upper(), 
         f"Общий уровень подготовки сотрудников МКДЦ. Средний балл {comparison}.", True),
        
        (Constants.TILE_X_LEFT, Constants.TILE_Y_BOTTOM, "ЭФФЕКТИВНОСТЬ ТЕСТИРОВАНИЯ", f"{passed_pct:.0f}%", 
         f"Сотрудников ({passed_count} из {total_participants} чел.) успешно преодолели порог в 70% баллов.", False),
        
        (Constants.TILE_X_RIGHT, Constants.TILE_Y_TOP, "АНАЛИЗ СЛОЖНЫХ ВОПРОСОВ", f"{min(q_success.values()) if q_success else 0:.0f}%", 
         f"Правильных ответов на самый сложный вопрос отчета:\n\"{worst_q[:55]}...\"", False),
        
        (Constants.TILE_X_RIGHT, Constants.TILE_Y_BOTTOM, "РЕЙТИНГ ПОДРАЗДЕЛЕНИЙ", f"{best_dep_pct:.0f}%", 
         f"Средний результат лучшего отделения МКДЦ:\n{best_dep}.", False)
    ]

    for x, y, tag, main_text, sub_text, is_accent in tiles_config:
        tile_shape = slide.shapes.add_shape(1, x, y, Constants.TILE_W, Constants.TILE_H)
        tile_shape.fill.solid()
        tile_shape.fill.fore_color.rgb = Constants.COLOR_TILE_BG
        
        tile_shape.line.color.rgb = Constants.COLOR_SECONDARY
        tile_shape.line.width = Constants.TILE_LINE_WIDTH
        
        if hasattr(tile_shape, "shadow") and tile_shape.shadow:
            tile_shape.shadow.inherit = False

        tx_box = slide.shapes.add_textbox(x + Inches(0.35), y + Inches(0.2), Constants.TILE_W - Inches(0.7), Constants.TILE_H - Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p_tag = tf.paragraphs[0]
        p_tag.text = tag
        p_tag.font.name = Constants.FONT_NAME
        p_tag.font.size = Constants.FONT_SIZE_TAG
        p_tag.font.bold = True
        p_tag.font.color.rgb = Constants.COLOR_MUTED_TAG

        p_main = tf.add_paragraph()
        p_main.text = main_text
        p_main.font.name = Constants.FONT_NAME
        p_main.font.size = Constants.FONT_SIZE_MAIN
        p_main.font.bold = True
        p_main.font.color.rgb = Constants.COLOR_ACCENT if is_accent else Constants.COLOR_PRIMARY
        p_main.space_before = Pt(1)

        p_sub = tf.add_paragraph()
        p_sub.text = sub_text
        p_sub.font.name = Constants.FONT_NAME
        p_sub.font.size = Constants.FONT_SIZE_SUB
        p_sub.font.color.rgb = Constants.COLOR_PRIMARY
        p_sub.space_before = Pt(4)
        p_sub.line_spacing = 1.2